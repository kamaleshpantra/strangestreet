#!/usr/bin/env python3
"""
Generate a PowerPoint deck: ML Intelligence System for Social Media (Strange Street).

Usage:
  pip install python-pptx
  python ml/build_ml_pipeline_presentation.py
  # Output: ml/StrangeStreet_ML_Intelligence_System.pptx
"""

from __future__ import annotations

import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
except ImportError:
    print("Install dependency: pip install python-pptx", file=sys.stderr)
    sys.exit(1)


def _set_slide_title(slide, text: str, subtitle: str | None = None):
    title = slide.shapes.title
    title.text = text
    tf = title.text_frame
    for p in tf.paragraphs:
        p.font.size = Pt(32)
        p.font.bold = True
        p.font.color.rgb = RGBColor(30, 30, 40)
    if subtitle and len(slide.placeholders) > 1:
        body = slide.placeholders[1]
        tf = body.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.text = subtitle
        p.font.size = Pt(18)
        p.font.color.rgb = RGBColor(80, 80, 100)


def _bullet_slide(prs, title: str, bullets: list[str]):
    layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.clear()
    for i, line in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = line
        p.level = 0
        p.font.size = Pt(20)
        p.space_after = Pt(10)
    return slide


def _two_column_bullets(prs, title: str, left_title: str, left: list[str], right_title: str, right: list[str]):
    """Title + manually placed text boxes for two columns."""
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Title
    box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(9), Inches(0.9))
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(28)
    p.font.bold = True

    def add_col(x, y, w, col_title, items):
        tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(5.5))
        tff = tb.text_frame
        tff.word_wrap = True
        pp = tff.paragraphs[0]
        pp.text = col_title
        pp.font.bold = True
        pp.font.size = Pt(22)
        pp.font.color.rgb = RGBColor(180, 40, 60)
        for it in items:
            p2 = tff.add_paragraph()
            p2.text = it
            p2.level = 0
            p2.font.size = Pt(16)
            p2.space_after = Pt(6)

    add_col(0.5, 1.3, 4.4, left_title, left)
    add_col(5.1, 1.3, 4.4, right_title, right)
    return slide


def build_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # --- Slide 1: Title ---
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    _set_slide_title(
        slide,
        "ML Intelligence System",
        "Social Media Application — Strange Street",
    )
    if slide.shapes.title.text_frame.paragraphs:
        slide.shapes.title.text_frame.paragraphs[0].font.size = Pt(36)
    sub = slide.placeholders[1].text_frame.paragraphs[0]
    sub.text = (
        "Batch intelligence layer: graph + collaborative filtering + semantic features\n"
        "Precomputed scores for fast feeds, discovery, zones, and safety"
    )
    sub.font.size = Pt(16)

    # --- Slide 2: What problem does it solve? ---
    _bullet_slide(
        prs,
        "What Problem Does It Solve?",
        [
            "Traditional social feeds favor people you already know — cold start for new interests & cities",
            "Users need: relevant posts, stranger discovery, communities (zones), and safe experiences",
            "Challenge: rich personalization is expensive if computed on every page load",
            "Solution: separate heavy ML into a pipeline; API serves precomputed scores (low latency)",
        ],
    )

    # --- Slide 3: Intelligence vs raw app ---
    _two_column_bullets(
        prs,
        "Application vs ML Intelligence Layer",
        "Core product (FastAPI + DB)",
        [
            "Auth, posts, comments, follows",
            "Zones, messages, connections",
            "Stores InteractionLog events",
            "Reads FeedScore, PeopleScore, ZoneScore",
        ],
        "ML intelligence (batch pipeline)",
        [
            "Graph: influence & communities",
            "Features: text / interest vectors",
            "Recommend: feed, people, zones",
            "Safety: toxicity heuristics",
            "Evaluation: offline metrics",
        ],
    )

    # --- Slide 4: Architecture ---
    _bullet_slide(
        prs,
        "High-Level Architecture",
        [
            "Data plane: PostgreSQL — users, posts, follows, interaction logs, ML cache tables",
            "Batch job: ml/run_pipeline.py (or cron on Render) runs on a schedule",
            "Flow: raw events → feature store → scoring tables → fast SELECT at request time",
            "Optional: APScheduler inside web process (daily) + external cron (hourly) for scale patterns",
        ],
    )

    # --- Slide 5: Data the system learns from ---
    _bullet_slide(
        prs,
        "Signals the System Uses",
        [
            "Social graph: who follows whom (directed edges)",
            "Implicit feedback: views, likes, comments, skips — weighted in InteractionLog",
            "Explicit interests: user ↔ interest tags",
            "Text: post content → embeddings (Sentence-BERT in production)",
            "Zones: memberships and posts-in-zone for community recommendations",
        ],
    )

    # --- Slide 6: Pipeline overview ---
    _bullet_slide(
        prs,
        "ML Pipeline — Seven Stages",
        [
            "1. Graph engine — PageRank, Louvain communities, friend-of-friend, sparse interest hints",
            "2. Feature engine — post & user semantic vectors; TruncatedSVD on interest matrix",
            "3. Feed recommender — TruncatedSVD on user×post interactions + recency + semantic re-rank",
            "4. People recommender — blend FoF, PageRank, community, Jaccard interests, topic cosine",
            "5. Zone recommender — interest + semantic centroid + member/community overlap + activity",
            "6. Safety — regex / keyword toxicity & spam scoring; flags + feed filtering",
            "7. Evaluation — Precision@K, catalog coverage, category diversity",
        ],
    )

    # --- Slide 7: Graph intelligence ---
    _bullet_slide(
        prs,
        "Graph Intelligence (Social Network Analysis)",
        [
            "PageRank: estimates influence in the follow graph (who is structurally central)",
            "Louvain: finds organic clusters — users who densely follow each other",
            "Friend-of-friend: 2-hop suggestions — “people near you” without direct follow",
            "Used for: boosting content, people suggestions, zone affinity via community overlap",
        ],
    )

    # --- Slide 8: Semantic & interest features ---
    _bullet_slide(
        prs,
        "Semantic & Interest Features",
        [
            "Post embedding: sentence transformer (e.g. all-MiniLM-L6-v2) → dense vector per post",
            "User profile vector: average of embeddings of posts they authored (writing style/topics)",
            "Interest SVD: binary user×interest matrix compressed — smooth similarity in tag space",
            "Why: combine “what you write” with “what tags you picked” for richer matching",
        ],
    )

    # --- Slide 9: Feed ML ---
    _bullet_slide(
        prs,
        "Smart Feed — Collaborative Filtering + Re-ranking",
        [
            "Build sparse matrix: users × posts, cell = sum of interaction weights",
            "TruncatedSVD (matrix factorization): latent “taste” factors for users and posts",
            "Score ≈ dot product of factors; blend with recency (fresh content)",
            "Stage 2: re-rank with topic similarity + graph signals; cap posts per author (diversity)",
            "Cold start: popular posts for users with no scores yet",
        ],
    )

    # --- Slide 10: Discovery ---
    _bullet_slide(
        prs,
        "Discovery — People & Zones",
        [
            "People: weighted score — FoF, target PageRank, same community, Jaccard on interests, cosine on embeddings",
            "Explainability: PeopleScore stores a JSON breakdown per candidate",
            "Zones: match user to community via interest overlap, semantic centroid of members, activity level",
            "Aligns product goal: “strangers first” via interests, not only existing friends",
        ],
    )

    # --- Slide 11: Safety ---
    _bullet_slide(
        prs,
        "Safety Layer (Moderation v1)",
        [
            "Heuristic toxicity & spam patterns (regex) → score 0–1; threshold flags posts",
            "Updates PostFeature.toxicity_score; high-toxic posts filtered from ML-ranked feed",
            "Author caps in feed scores reduce dominance by one creator",
            "Roadmap note: can swap heuristics for a trained classifier (e.g. Jigsaw-style)",
        ],
    )

    # --- Slide 12: Evaluation ---
    _bullet_slide(
        prs,
        "Offline Evaluation (Sanity Checks)",
        [
            "Precision@K: top-K feed vs posts user actually liked/commented (historical proxy)",
            "Coverage: fraction of catalog that appears in recommendations (exploration)",
            "Diversity: spread of categories in recommended sets",
            "Use: compare pipeline versions — not a substitute for A/B tests with real users",
        ],
    )

    # --- Slide 13: Deployment ---
    _bullet_slide(
        prs,
        "Deployment & Operations",
        [
            "Production stack: FastAPI, PostgreSQL (e.g. Neon), optional Cloudinary for media",
            "Gunicorn + Uvicorn workers for the API; ML as separate cron/worker job recommended at scale",
            "Render blueprint: web service + scheduled job running ml/run_pipeline.py",
            "Disk / DB persistence for JSON-like score tables — same pattern as simulation, but durable",
        ],
    )

    # --- Slide 14: Simulation tool ---
    _bullet_slide(
        prs,
        "In-Memory Simulation (ml/simulate_pipeline_demo.py)",
        [
            "Generates toy or large synthetic social data — no database required",
            "Prints embeddings, factors, scores, and metrics for learning and debugging",
            "Flags: --large, --users/--posts/--interactions, --verbose for full trace",
            "Educational mirror of production logic — not wired to live traffic",
        ],
    )

    # --- Slide 15: Takeaways ---
    _bullet_slide(
        prs,
        "Key Takeaways",
        [
            "Treat ML as a subsystem: ingest events → batch compute → serve scores",
            "Blend graph structure, behavior (CF), and semantics for social products",
            "Precompute for speed; explain partial scores where possible (people breakdown)",
            "Safety and diversity are first-class outputs, not afterthoughts",
            "Validate offline, then validate online with real users",
        ],
    )

    return prs


def main():
    out_dir = os.path.dirname(os.path.abspath(__file__))
    out_path = os.path.join(out_dir, "StrangeStreet_ML_Intelligence_System.pptx")
    prs = build_deck()
    prs.save(out_path)
    print(f"Saved: {out_path}")


if __name__ == "__main__":
    main()
