"""
Build a 10-slide PPT for the Strange Street ML Pipeline presentation.
Run: python build_ml_ppt.py
Output: ML_Pipeline_Presentation.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Theme Colors ──────────────────────────────────────────────────────────
BG_DARK      = RGBColor(0x0F, 0x0F, 0x1A)   # near-black blue
BG_CARD      = RGBColor(0x1A, 0x1C, 0x2E)   # card background
ACCENT_BLUE  = RGBColor(0x60, 0xA5, 0xFA)   # sky blue
ACCENT_PURPLE= RGBColor(0xA7, 0x8B, 0xFA)   # lavender
ACCENT_GREEN = RGBColor(0x34, 0xD3, 0x99)   # emerald
ACCENT_ORANGE= RGBColor(0xFB, 0x92, 0x3C)   # orange
ACCENT_PINK  = RGBColor(0xF4, 0x72, 0xB6)   # pink
ACCENT_RED   = RGBColor(0xF8, 0x71, 0x71)   # coral red
TEXT_WHITE    = RGBColor(0xF1, 0xF5, 0xF9)   # off-white
TEXT_GRAY     = RGBColor(0x94, 0xA3, 0xB8)   # muted gray
TEXT_DIM      = RGBColor(0x64, 0x74, 0x8B)   # dimmer gray


prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)

SLIDE_W = prs.slide_width
SLIDE_H = prs.slide_height


def set_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_shape(slide, left, top, w, h, fill_color, corner_radius=Inches(0.15)):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    shape.shadow.inherit = False
    # Set corner radius via adjustments
    if shape.adjustments and len(shape.adjustments) > 0:
        shape.adjustments[0] = 0.04
    return shape


def add_text(slide, left, top, w, h, text, font_size=18, color=TEXT_WHITE,
             bold=False, align=PP_ALIGN.LEFT, font_name="Segoe UI"):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox


def add_multiline(slide, left, top, w, h, lines, font_size=16,
                  color=TEXT_WHITE, spacing=1.2, font_name="Segoe UI",
                  bold_first=False):
    """lines = list of strings. Each becomes a paragraph."""
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.font.name = font_name
        p.space_after = Pt(font_size * spacing * 0.4)
        if bold_first and i == 0:
            p.font.bold = True
    return txBox


def add_bullet_card(slide, left, top, w, h, title, bullets, accent, 
                    title_size=20, bullet_size=15):
    """Card with colored title and bullet points."""
    card = add_shape(slide, left, top, w, h, BG_CARD)
    # Accent bar on top
    bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, w, Inches(0.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = accent
    bar.line.fill.background()
    # Title
    add_text(slide, left + Inches(0.25), top + Inches(0.15), w - Inches(0.5), Inches(0.5),
             title, font_size=title_size, color=accent, bold=True)
    # Bullets
    y = top + Inches(0.65)
    for b in bullets:
        add_text(slide, left + Inches(0.35), y, w - Inches(0.6), Inches(0.4),
                 f"→  {b}", font_size=bullet_size, color=TEXT_GRAY)
        y += Inches(0.38)
    return card


def add_metric_box(slide, left, top, w, h, label, value, accent):
    card = add_shape(slide, left, top, w, h, BG_CARD)
    add_text(slide, left, top + Inches(0.15), w, Inches(0.6),
             value, font_size=32, color=accent, bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, left, top + Inches(0.65), w, Inches(0.4),
             label, font_size=12, color=TEXT_DIM, align=PP_ALIGN.CENTER)
    return card


def slide_number(slide, num):
    add_text(slide, SLIDE_W - Inches(1), SLIDE_H - Inches(0.5),
             Inches(0.8), Inches(0.3), f"{num}/10",
             font_size=11, color=TEXT_DIM, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])  # blank
set_bg(sl)

# Decorative accent line
accent_line = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE,
    Inches(1.2), Inches(1.6), Inches(4), Inches(0.06))
accent_line.fill.solid()
accent_line.fill.fore_color.rgb = ACCENT_BLUE
accent_line.line.fill.background()

add_text(sl, Inches(1.2), Inches(2.0), Inches(10), Inches(1.2),
         "Strange Street", font_size=52, color=ACCENT_BLUE, bold=True,
         font_name="Segoe UI Semibold")

add_text(sl, Inches(1.2), Inches(3.2), Inches(10), Inches(0.8),
         "ML Intelligence Pipeline", font_size=36, color=TEXT_WHITE, bold=False,
         font_name="Segoe UI Light")

add_text(sl, Inches(1.2), Inches(4.3), Inches(8), Inches(0.5),
         "How machine learning powers personalized feeds, stranger discovery, and content safety",
         font_size=16, color=TEXT_GRAY)

add_text(sl, Inches(1.2), Inches(6.2), Inches(6), Inches(0.4),
         "Interest-Based Social Networking Platform  ·  FastAPI + scikit-learn + NetworkX",
         font_size=13, color=TEXT_DIM)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — The Big Picture
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
slide_number(sl, 2)

add_text(sl, Inches(0.8), Inches(0.4), Inches(8), Inches(0.7),
         "What Does the Pipeline Do?", font_size=34, color=TEXT_WHITE, bold=True)

add_text(sl, Inches(0.8), Inches(1.1), Inches(10), Inches(0.5),
         "Runs daily at 2:00 AM  —  transforms raw user data into personalized experiences",
         font_size=15, color=TEXT_GRAY)

# 7 step boxes in a flow
steps = [
    ("1", "Graph\nEngine",    ACCENT_BLUE),
    ("2", "Feature\nEngine",  ACCENT_PURPLE),
    ("3", "Feed\nRecommender",ACCENT_GREEN),
    ("4", "People\nRecommender",ACCENT_ORANGE),
    ("5", "Zone\nRecommender",ACCENT_PINK),
    ("6", "Safety\nModule",   ACCENT_RED),
    ("7", "Evaluation",       ACCENT_BLUE),
]

start_x = Inches(0.5)
y = Inches(2.3)
box_w = Inches(1.55)
box_h = Inches(1.4)
gap = Inches(0.2)

for i, (num, label, col) in enumerate(steps):
    x = start_x + i * (box_w + gap)
    card = add_shape(sl, x, y, box_w, box_h, BG_CARD)
    
    # Step number circle
    circle = sl.shapes.add_shape(MSO_SHAPE.OVAL,
        x + Inches(0.55), y + Inches(0.12), Inches(0.45), Inches(0.45))
    circle.fill.solid()
    circle.fill.fore_color.rgb = col
    circle.line.fill.background()
    add_text(sl, x + Inches(0.55), y + Inches(0.14), Inches(0.45), Inches(0.45),
             num, font_size=18, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)
    
    add_text(sl, x + Inches(0.05), y + Inches(0.65), box_w - Inches(0.1), Inches(0.8),
             label, font_size=13, color=TEXT_GRAY, align=PP_ALIGN.CENTER)
    
    # Arrow between boxes
    if i < 6:
        add_text(sl, x + box_w, y + Inches(0.45), Inches(0.2), Inches(0.3),
                 "›", font_size=22, color=TEXT_DIM, align=PP_ALIGN.CENTER)

# Three output cards at bottom
outputs = [
    ("📱  Personalized Feed", "Each user sees posts ranked by\nhow likely they are to engage\nwith them", ACCENT_GREEN),
    ("👤  Stranger Discovery", "Recommends strangers based on\nshared interests, graph proximity\nand content similarity", ACCENT_ORANGE),
    ("🛡️  Content Safety", "Automatically flags toxic posts\nand enforces diversity in feeds", ACCENT_RED),
]

y2 = Inches(4.4)
card_w = Inches(3.8)
for i, (title, desc, col) in enumerate(outputs):
    x = Inches(0.6) + i * (card_w + Inches(0.3))
    card = add_shape(sl, x, y2, card_w, Inches(2.2), BG_CARD)
    bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y2, Inches(0.06), Inches(2.2))
    bar.fill.solid()
    bar.fill.fore_color.rgb = col
    bar.line.fill.background()
    add_text(sl, x + Inches(0.25), y2 + Inches(0.2), card_w - Inches(0.4), Inches(0.4),
             title, font_size=17, color=col, bold=True)
    add_text(sl, x + Inches(0.25), y2 + Inches(0.65), card_w - Inches(0.4), Inches(1.4),
             desc, font_size=13, color=TEXT_DIM)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Graph Engine
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
slide_number(sl, 3)

add_text(sl, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
         "Step 1 — Graph Engine", font_size=34, color=ACCENT_BLUE, bold=True)
add_text(sl, Inches(0.8), Inches(1.05), Inches(10), Inches(0.4),
         "Analyzing the social network structure using graph theory",
         font_size=15, color=TEXT_GRAY)

# Four algorithm cards
algos = [
    ("PageRank", 
     ["Google's algorithm adapted for users",
      "Measures influence based on\nwho follows whom",
      "Users followed by influential\npeople score higher"],
     ACCENT_BLUE, "NetworkX\nnx.pagerank()"),
    
    ("Louvain Communities",
     ["Automatically discovers groups\nof tightly-connected users",
      "Maximizes modularity — internal\nconnections vs external",
      "No need to specify number\nof communities upfront"],
     ACCENT_PURPLE, "python-louvain\nbest_partition()"),
    
    ("Friend-of-Friend",
     ["Finds people 2 hops away\nin the social graph",
      "\"You might know them through\na mutual connection\"",
      "Capped at 50 FoF per user"],
     ACCENT_GREEN, "NetworkX\nDiGraph traversal"),
    
    ("Label Propagation",
     ["Infers interests for users with\nsparse profiles (< 3 interests)",
      "Looks at what neighbors like\nand suggests the most common",
      "Fills gaps so recommendations\nwork for everyone"],
     ACCENT_ORANGE, "Custom\nneighbor counting"),
]

card_w = Inches(2.85)
card_h = Inches(4.4)
start_x = Inches(0.5)
y = Inches(1.7)

for i, (title, bullets, col, lib) in enumerate(algos):
    x = start_x + i * (card_w + Inches(0.2))
    card = add_shape(sl, x, y, card_w, card_h, BG_CARD)
    
    # Color bar
    bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = col
    bar.line.fill.background()
    
    add_text(sl, x + Inches(0.2), y + Inches(0.15), card_w - Inches(0.4), Inches(0.4),
             title, font_size=19, color=col, bold=True)
    
    by = y + Inches(0.65)
    for b in bullets:
        add_text(sl, x + Inches(0.2), by, card_w - Inches(0.4), Inches(0.7),
                 b, font_size=12, color=TEXT_GRAY)
        by += Inches(0.72)
    
    # Library tag
    add_text(sl, x + Inches(0.15), y + card_h - Inches(0.55),
             card_w - Inches(0.3), Inches(0.4),
             lib, font_size=10, color=TEXT_DIM)

# Output note
add_text(sl, Inches(0.8), Inches(6.4), Inches(11), Inches(0.4),
         "Output  ▸  Per-user features: pagerank score, community ID, degree, FoF set, inferred interests",
         font_size=13, color=TEXT_DIM)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Feature Engine
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
slide_number(sl, 4)

add_text(sl, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
         "Step 2 — Feature Engine", font_size=34, color=ACCENT_PURPLE, bold=True)
add_text(sl, Inches(0.8), Inches(1.05), Inches(10), Inches(0.4),
         "Converting raw data into numerical vectors that algorithms can process",
         font_size=15, color=TEXT_GRAY)

# Three feature types
features = [
    ("Text Features", "Sentence-BERT", ACCENT_BLUE,
     ["Model: all-MiniLM-L6-v2",
      "Input: post text content",
      "Output: 384-dimensional vector",
      "Similar posts → similar vectors",
      "User profile = average of their\npost vectors (centroid)"]),
    
    ("Interest Features", "Truncated SVD", ACCENT_PURPLE,
     ["Input: user × interest binary matrix\n(300+ interests, mostly zeros)",
      "SVD compresses to 50 dimensions",
      "Captures hidden patterns:\nusers who like Python + ML\nare close in this space",
      "Explained variance tracks\ninformation preserved"]),
    
    ("Behavioral Features", "Statistical", ACCENT_GREEN,
     ["Engagement Rate:\n(likes + comments) / total actions",
      "Activity Level:\nactions in last 30 days / 30",
      "Quantifies how active and\ninteractive each user is",
      "Used for quality weighting\nin feed re-ranking"]),
]

card_w = Inches(3.8)
card_h = Inches(4.6)
start_x = Inches(0.5)
y = Inches(1.7)

for i, (title, method, col, bullets) in enumerate(features):
    x = start_x + i * (card_w + Inches(0.25))
    card = add_shape(sl, x, y, card_w, card_h, BG_CARD)
    
    # Header with method tag
    add_text(sl, x + Inches(0.25), y + Inches(0.15), card_w - Inches(0.5), Inches(0.35),
             title, font_size=20, color=col, bold=True)
    
    # Method tag
    tag = add_shape(sl, x + Inches(0.25), y + Inches(0.55),
                    Inches(1.8), Inches(0.3), col)
    add_text(sl, x + Inches(0.25), y + Inches(0.55), Inches(1.8), Inches(0.3),
             method, font_size=11, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)
    
    by = y + Inches(1.05)
    for b in bullets:
        add_text(sl, x + Inches(0.25), by, card_w - Inches(0.5), Inches(0.8),
                 b, font_size=12, color=TEXT_GRAY)
        by += Inches(0.7)

# Output
add_text(sl, Inches(0.8), Inches(6.6), Inches(11), Inches(0.4),
         "Output  ▸  UserFeature table (pagerank, embeddings, engagement) + PostFeature table (topic vectors, toxicity)",
         font_size=13, color=TEXT_DIM)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Feed Recommender
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
slide_number(sl, 5)

add_text(sl, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
         "Step 3 — Feed Recommender", font_size=34, color=ACCENT_GREEN, bold=True)
add_text(sl, Inches(0.8), Inches(1.05), Inches(10), Inches(0.4),
         "Collaborative filtering — \"users who liked similar posts will like similar future posts\"",
         font_size=15, color=TEXT_GRAY)

# Stage boxes - horizontal flow
stages = [
    ("Build Matrix", "User × Post interaction matrix\nWeights: view=0.1, like=1.0\ncomment=2.0, share=3.0", "1"),
    ("SVD Factorize", "Decompose into user_factors\n× post_factors (50 latent dims)\nDiscovers hidden preferences", "2"),
    ("Recency Boost", "Exponential decay: e^(-age/30)\nFresh posts ≈ 1.0 boost\n30-day posts ≈ 0.37 boost", "3"),
    ("Re-Rank", "Blend with topic similarity\nand PageRank features\n80% SVD + 20% features", "4"),
]

stage_w = Inches(2.7)
stage_h = Inches(2.6)
sx = Inches(0.5)
sy = Inches(1.8)

for i, (title, desc, num) in enumerate(stages):
    x = sx + i * (stage_w + Inches(0.35))
    card = add_shape(sl, x, sy, stage_w, stage_h, BG_CARD)
    
    # Number
    add_text(sl, x + Inches(0.15), sy + Inches(0.1), Inches(0.4), Inches(0.4),
             num, font_size=26, color=ACCENT_GREEN, bold=True)
    
    add_text(sl, x + Inches(0.5), sy + Inches(0.12), stage_w - Inches(0.6), Inches(0.35),
             title, font_size=17, color=TEXT_WHITE, bold=True)
    
    add_text(sl, x + Inches(0.2), sy + Inches(0.65), stage_w - Inches(0.4), Inches(1.8),
             desc, font_size=12, color=TEXT_GRAY)
    
    if i < 3:
        add_text(sl, x + stage_w + Inches(0.05), sy + Inches(0.9),
                 Inches(0.25), Inches(0.4), "→", font_size=20, color=TEXT_DIM)

# Bottom: special handling cards
bottom_y = Inches(4.8)
special = [
    ("Author Diversity Cap", "Max 3 posts per author in any\nuser's feed — prevents domination", ACCENT_ORANGE),
    ("Cold Start Fallback", "New users with no history get\npopularity-ranked posts (most liked)", ACCENT_PINK),
    ("Final Formula", "score = 0.7 × SVD + 0.3 × recency\nthen re-ranked with topic + PageRank", ACCENT_BLUE),
]

for i, (title, desc, col) in enumerate(special):
    x = Inches(0.5) + i * (Inches(4.0) + Inches(0.2))
    card = add_shape(sl, x, bottom_y, Inches(4.0), Inches(1.8), BG_CARD)
    bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, bottom_y, Inches(0.05), Inches(1.8))
    bar.fill.solid()
    bar.fill.fore_color.rgb = col
    bar.line.fill.background()
    add_text(sl, x + Inches(0.25), bottom_y + Inches(0.15), Inches(3.6), Inches(0.35),
             title, font_size=15, color=col, bold=True)
    add_text(sl, x + Inches(0.25), bottom_y + Inches(0.55), Inches(3.6), Inches(1.1),
             desc, font_size=12, color=TEXT_DIM)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — People Recommender
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
slide_number(sl, 6)

add_text(sl, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
         "Step 4 — People Recommender", font_size=34, color=ACCENT_ORANGE, bold=True)
add_text(sl, Inches(0.8), Inches(1.05), Inches(10), Inches(0.4),
         "5-signal scoring system to recommend the best strangers to connect with",
         font_size=15, color=TEXT_GRAY)

# 5 signal bars
signals = [
    ("Interest Jaccard", "30%", "|A ∩ B| / |A ∪ B|", ACCENT_ORANGE, 0.30),
    ("Friend-of-Friend", "25%", "Is target in your FoF set?", ACCENT_BLUE, 0.25),
    ("Topic Similarity", "20%", "Cosine similarity of BERT vectors", ACCENT_PURPLE, 0.20),
    ("Same Community", "15%", "Both in same Louvain cluster?", ACCENT_GREEN, 0.15),
    ("PageRank", "10%", "Target's influence score", ACCENT_PINK, 0.10),
]

sy = Inches(1.9)
bar_max_w = Inches(7.0)

for i, (name, pct, desc, col, frac) in enumerate(signals):
    y = sy + i * Inches(0.95)
    
    # Label
    add_text(sl, Inches(0.8), y, Inches(2.5), Inches(0.35),
             name, font_size=16, color=TEXT_WHITE, bold=True)
    
    # Bar background
    bg_bar = add_shape(sl, Inches(3.5), y + Inches(0.05), bar_max_w, Inches(0.35), BG_CARD)
    
    # Filled bar  
    fill_bar = add_shape(sl, Inches(3.5), y + Inches(0.05),
                         Inches(bar_max_w.inches * frac / 0.30), Inches(0.35), col)
    
    # Percentage
    add_text(sl, Inches(3.5) + Inches(bar_max_w.inches * frac / 0.30) + Inches(0.15),
             y + Inches(0.02), Inches(0.6), Inches(0.35),
             pct, font_size=14, color=col, bold=True)
    
    # Description
    add_text(sl, Inches(3.5), y + Inches(0.42), bar_max_w, Inches(0.3),
             desc, font_size=11, color=TEXT_DIM)

# Example box
ex_y = Inches(6.6)
add_text(sl, Inches(0.8), Inches(6.5), Inches(12), Inches(0.4),
         "Example  ▸  score = 0.25×1.0 + 0.10×0.45 + 0.15×1.0 + 0.30×0.50 + 0.20×0.68 = 0.731  →  strong match",
         font_size=13, color=TEXT_DIM)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Zone Recommender
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
slide_number(sl, 7)

add_text(sl, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
         "Step 5 — Zone Recommender", font_size=34, color=ACCENT_PINK, bold=True)
add_text(sl, Inches(0.8), Inches(1.05), Inches(10), Inches(0.4),
         "Recommends communities (Zones) you haven't joined but would enjoy",
         font_size=15, color=TEXT_GRAY)

# 5 signals as cards in a 3+2 grid
z_signals = [
    ("Semantic Overlap", "25%",
     "BERT cosine similarity between your\ntopic profile and the zone's average\nmember topic profile",
     ACCENT_PURPLE),
    ("Binary Interest", "20%",
     "Jaccard overlap between your interests\nand the zone members' aggregated\ninterest set",
     ACCENT_ORANGE),
    ("Activity Level", "20%",
     "Zone's post count relative to the\nmost active zone on the platform\nPrefers active communities",
     ACCENT_GREEN),
    ("Member Overlap", "20%",
     "How many people you follow are\nalready in this zone?\nMax score at 3+ friends in zone",
     ACCENT_BLUE),
    ("Community Match", "15%",
     "What fraction of zone members\nare in your Louvain community?\nHigher = better cultural fit",
     ACCENT_PINK),
]

card_w = Inches(3.8)
card_h = Inches(2.0)

for i, (title, pct, desc, col) in enumerate(z_signals):
    row = 0 if i < 3 else 1
    col_idx = i if i < 3 else i - 3
    
    x = Inches(0.5) + col_idx * (card_w + Inches(0.25))
    y = Inches(1.8) + row * (card_h + Inches(0.3))
    
    if row == 1:
        x = Inches(0.5) + col_idx * (card_w + Inches(0.25)) + Inches(2.0)
    
    card = add_shape(sl, x, y, card_w, card_h, BG_CARD)
    
    # Percentage badge
    badge = add_shape(sl, x + card_w - Inches(0.8), y + Inches(0.1),
                      Inches(0.65), Inches(0.3), col)
    add_text(sl, x + card_w - Inches(0.8), y + Inches(0.1),
             Inches(0.65), Inches(0.3), pct,
             font_size=12, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)
    
    add_text(sl, x + Inches(0.2), y + Inches(0.12), card_w - Inches(1.2), Inches(0.35),
             title, font_size=17, color=col, bold=True)
    
    add_text(sl, x + Inches(0.2), y + Inches(0.55), card_w - Inches(0.4), Inches(1.3),
             desc, font_size=12, color=TEXT_DIM)

# Note
add_text(sl, Inches(0.8), Inches(6.6), Inches(11), Inches(0.4),
         "Output  ▸  Top 10 zone recommendations per user stored in ZoneScore table",
         font_size=13, color=TEXT_DIM)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Safety Module
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
slide_number(sl, 8)

add_text(sl, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
         "Step 6 — Safety Module", font_size=34, color=ACCENT_RED, bold=True)
add_text(sl, Inches(0.8), Inches(1.05), Inches(10), Inches(0.4),
         "Protecting users from harmful content before it reaches their feeds",
         font_size=15, color=TEXT_GRAY)

# Two main cards side by side
# Left: Toxicity Detection
add_bullet_card(sl, Inches(0.5), Inches(1.8), Inches(5.8), Inches(4.5),
    "Toxicity Detection", [
        "Regex pattern matching against known harmful patterns",
        "Toxic: threats, hate speech, slurs, targeted insults",
        "Spam: promotions, repeated characters, multi-URL posts",
        "Score = pattern hits / threshold — range 0.0 to 1.0",
        "If score ≥ 0.3 → post is auto-flagged and hidden",
        "ContentFlag record created for audit trail",
        "Future: upgrade to Jigsaw trained classifier (v2)",
    ], ACCENT_RED, title_size=22, bullet_size=14)

# Right: Author Diversity
add_bullet_card(sl, Inches(6.8), Inches(1.8), Inches(5.8), Inches(4.5),
    "Feed Diversity Enforcement", [
        "Max 3 posts per author in any user's feed",
        "Prevents popular authors from dominating",
        "Sorted by score — keeps the BEST 3, drops rest",
        "Runs after all scoring is complete",
        "Ensures variety and fair exposure for all creators",
        "Applies to both SVD-scored and fallback feeds",
    ], ACCENT_ORANGE, title_size=22, bullet_size=14)

# Scoring formula box
add_shape(sl, Inches(0.5), Inches(6.5), Inches(12.3), Inches(0.7), BG_CARD)
add_text(sl, Inches(0.8), Inches(6.55), Inches(11.5), Inches(0.6),
         "Scoring  ▸  toxicity = min( hits / (total_patterns × 0.3),  1.0 )    |    threshold for auto-flag = 0.3",
         font_size=14, color=TEXT_DIM)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Evaluation Metrics
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
slide_number(sl, 9)

add_text(sl, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
         "Step 7 — Evaluation Metrics", font_size=34, color=ACCENT_BLUE, bold=True)
add_text(sl, Inches(0.8), Inches(1.05), Inches(10), Inches(0.4),
         "Measuring pipeline quality — \"Is our ML actually helping users?\"",
         font_size=15, color=TEXT_GRAY)

# Three metric cards
metrics = [
    ("Precision@K", ACCENT_GREEN,
     "Of the top-K recommended posts,\nhow many did the user actually\nlike or comment on?",
     "hits in top K / K",
     "Target: 0.3 – 0.7"),
    
    ("Catalog Coverage", ACCENT_BLUE,
     "What fraction of ALL posts appear\nin at least one user's feed?\nHigher = fairer for small creators",
     "unique recommended / total posts",
     "Target: 60% – 90%"),
    
    ("Feed Diversity", ACCENT_PURPLE,
     "Average number of distinct post\ncategories per user's feed\nHigher = less filter bubble",
     "avg unique categories per feed",
     "Target: 3 – 6 categories"),
]

card_w = Inches(3.8)
card_h = Inches(4.0)
start_x = Inches(0.5)
y = Inches(1.8)

for i, (title, col, desc, formula, target) in enumerate(metrics):
    x = start_x + i * (card_w + Inches(0.3))
    card = add_shape(sl, x, y, card_w, card_h, BG_CARD)
    
    bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, card_w, Inches(0.05))
    bar.fill.solid()
    bar.fill.fore_color.rgb = col
    bar.line.fill.background()
    
    add_text(sl, x + Inches(0.25), y + Inches(0.2), card_w - Inches(0.5), Inches(0.4),
             title, font_size=22, color=col, bold=True)
    
    add_text(sl, x + Inches(0.25), y + Inches(0.7), card_w - Inches(0.5), Inches(1.5),
             desc, font_size=14, color=TEXT_GRAY)
    
    # Formula tag
    add_shape(sl, x + Inches(0.2), y + Inches(2.1), card_w - Inches(0.4), Inches(0.4), 
              RGBColor(0x12, 0x14, 0x24))
    add_text(sl, x + Inches(0.2), y + Inches(2.15), card_w - Inches(0.4), Inches(0.35),
             formula, font_size=12, color=col, align=PP_ALIGN.CENTER)
    
    # Target
    add_text(sl, x + Inches(0.25), y + Inches(2.7), card_w - Inches(0.5), Inches(0.4),
             target, font_size=15, color=TEXT_WHITE, bold=True)

# Also: Pipeline Summary
add_shape(sl, Inches(0.5), Inches(6.1), Inches(12.3), Inches(1.0), BG_CARD)
add_text(sl, Inches(0.8), Inches(6.15), Inches(3), Inches(0.4),
         "Pipeline Health Check", font_size=16, color=TEXT_WHITE, bold=True)
add_text(sl, Inches(0.8), Inches(6.55), Inches(11), Inches(0.4),
         "Also tracks: total feed scores generated  ·  people scores  ·  zone scores  ·  content flags count",
         font_size=13, color=TEXT_DIM)


# ══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Tech Stack & Summary
# ══════════════════════════════════════════════════════════════════════════
sl = prs.slides.add_slide(prs.slide_layouts[6])
set_bg(sl)
slide_number(sl, 10)

add_text(sl, Inches(0.8), Inches(0.4), Inches(10), Inches(0.7),
         "Architecture & Tech Stack", font_size=34, color=TEXT_WHITE, bold=True)

# Algorithm table
headers = ["Algorithm", "Library", "Purpose"]
rows = [
    ["PageRank",             "NetworkX",             "User influence scoring"],
    ["Louvain Detection",    "python-louvain",       "Community clustering"],
    ["Sentence-BERT",        "sentence-transformers", "Semantic text embeddings (384d)"],
    ["Truncated SVD",        "scikit-learn",         "Matrix factorization + dim reduction"],
    ["Jaccard Similarity",   "Custom",               "Interest set overlap"],
    ["Cosine Similarity",    "NumPy",                "Vector similarity (topics, search)"],
    ["Exponential Decay",    "NumPy",                "Recency boost for fresh content"],
    ["Regex Matching",       "Python re",            "Toxicity & spam detection"],
]

table_x = Inches(0.5)
table_y = Inches(1.3)
col_widths = [Inches(2.4), Inches(2.6), Inches(4.5)]

# Header row
for j, (header, cw) in enumerate(zip(headers, col_widths)):
    x = table_x + sum(w.inches for w in col_widths[:j]) * 914400 / 914400
    x_calc = table_x
    for k in range(j):
        x_calc += col_widths[k]
    
    hdr_shape = add_shape(sl, x_calc, table_y, cw, Inches(0.45), ACCENT_BLUE)
    add_text(sl, x_calc + Inches(0.15), table_y + Inches(0.05), cw - Inches(0.3), Inches(0.35),
             header, font_size=13, color=BG_DARK, bold=True)

# Data rows
for i, row_data in enumerate(rows):
    ry = table_y + Inches(0.5) + i * Inches(0.42)
    bg = BG_CARD if i % 2 == 0 else RGBColor(0x14, 0x16, 0x26)
    
    for j, (cell, cw) in enumerate(zip(row_data, col_widths)):
        cx = table_x
        for k in range(j):
            cx += col_widths[k]
        
        cell_shape = add_shape(sl, cx, ry, cw, Inches(0.4), bg)
        col = ACCENT_BLUE if j == 0 else (ACCENT_PURPLE if j == 1 else TEXT_GRAY)
        bld = j == 0
        add_text(sl, cx + Inches(0.15), ry + Inches(0.05), cw - Inches(0.3), Inches(0.3),
                 cell, font_size=12, color=col, bold=bld)

# Key takeaways
ky = Inches(5.1)
takeaways = [
    ("Lightweight ML", "Runs on CPU, no GPU needed — fits free hosting tiers", ACCENT_GREEN),
    ("Graceful Fallbacks", "Every ML feature has a non-ML fallback for cold starts", ACCENT_ORANGE),
    ("Privacy First", "All processing is server-side — no user data leaves the platform", ACCENT_PURPLE),
]

for i, (title, desc, col) in enumerate(takeaways):
    x = Inches(0.5) + i * (Inches(4.1) + Inches(0.15))
    card = add_shape(sl, x, ky, Inches(4.1), Inches(1.6), BG_CARD)
    bar = sl.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, ky, Inches(0.05), Inches(1.6))
    bar.fill.solid()
    bar.fill.fore_color.rgb = col
    bar.line.fill.background()
    add_text(sl, x + Inches(0.25), ky + Inches(0.15), Inches(3.6), Inches(0.35),
             title, font_size=17, color=col, bold=True)
    add_text(sl, x + Inches(0.25), ky + Inches(0.55), Inches(3.6), Inches(0.9),
             desc, font_size=13, color=TEXT_DIM)

# Thank you
add_text(sl, Inches(0.8), Inches(6.9), Inches(11), Inches(0.4),
         "Strange Street  —  An interest-based social platform powered by 8 ML algorithms across a 7-step daily pipeline",
         font_size=13, color=TEXT_DIM)


# ══════════════════════════════════════════════════════════════════════════
# SAVE
# ══════════════════════════════════════════════════════════════════════════
output_path = os.path.join(os.path.dirname(__file__), "ML_Pipeline_Presentation.pptx")
prs.save(output_path)
print(f"\n✅ Presentation saved to: {output_path}")
print(f"   10 slides | Dark theme | Widescreen 16:9")
